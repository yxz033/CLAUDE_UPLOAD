import os
from pathlib import Path
import sys
import subprocess
import warnings
from PIL import Image
import multiprocessing as mp
from multiprocessing import Pool, cpu_count
import cv2
import numpy as np
import time
Image.MAX_IMAGE_PIXELS = None  # 禁用图片大小限制警告

# 添加用户目录的 site-packages 到 Python 路径
user_site_packages = os.path.expanduser("~/.local/lib/python3/site-packages")
if os.path.exists(user_site_packages) and user_site_packages not in sys.path:
    sys.path.append(user_site_packages)

# Windows 用户目录
windows_site_packages = os.path.expanduser("~/AppData/Roaming/Python/Python312/site-packages")
if os.path.exists(windows_site_packages) and windows_site_packages not in sys.path:
    sys.path.append(windows_site_packages)

def install_package(package_name):
    """安装Python包"""
    print(f"正在安装 {package_name}...")
    try:
        # 先尝试使用清华源安装
        try:
            subprocess.check_call([
                sys.executable, "-m", "pip", "install",
                "-i", "https://pypi.tuna.tsinghua.edu.cn/simple",
                package_name
            ])
            print(f"✓ {package_name} 安装成功")
            return True
        except:
            # 如果清华源失败,使用默认源
            print(f"使用清华源安装失败,尝试默认源...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
            print(f"✓ {package_name} 安装成功")
            return True
    except subprocess.CalledProcessError as e:
        print(f"× {package_name} 安装失败: {str(e)}")
        return False
    except Exception as e:
        print(f"× 安装过程出错: {str(e)}")
        return False

# 定义所需的依赖
dependencies = {
    "tqdm": "tqdm",
    "python-docx": "docx",
    "PyPDF2": "PyPDF2",
    "python-pptx": "pptx",
    "openpyxl": "openpyxl",
    "pywin32": "win32com",
    "Pillow": "PIL",
    "pytesseract": "pytesseract",
    "PyMuPDF": "fitz",
    "pdf2image": "pdf2image",
    "paddleocr": "paddleocr"  # 添加PaddleOCR
}

# 检查并安装依赖
print("检查依赖...")
for package, module in dependencies.items():
    try:
        __import__(module)
        print(f"✓ {package} 已安装")
    except ImportError:
        print(f"× {package} 未安装,开始安装...")
        if not install_package(package):
            print(f"! 请手动安装 {package}")
            sys.exit(1)

# 现在可以安全地导入所需的模块
try:
    from tqdm import tqdm
    from docx import Document
    from PyPDF2 import PdfReader
    from pptx import Presentation
    from openpyxl import load_workbook
    import win32com.client
    import re
    import io
    from PIL import Image
    import pytesseract
    import fitz
    from pdf2image import convert_from_path
    from paddleocr import PaddleOCR  # 添加PaddleOCR导入
except ImportError as e:
    print(f"导入模块时出错: {str(e)}")
    print("请确保所有依赖都已正确安装")
    sys.exit(1)

# 过滤PDF处理的警告
warnings.filterwarnings('ignore', category=UserWarning, module='PyPDF2')

# 设置Tesseract路径
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def get_optimal_processes():
    """获取最优进程数
    返回CPU核心数-1，确保系统留有余力
    最小返回1，最大返回CPU核心数
    """
    cpu_cores = cpu_count()
    return max(1, min(cpu_cores - 1, cpu_cores))

def check_dependencies():
    """检查所有必要的依赖是否已安装,如果没有则自动安装"""
    missing_deps = []
    installed_deps = []
    manual_deps = []  # 需要手动安装的依赖
    
    # 依赖列表
    dependencies = {
        "python-docx": "docx",
        "PyPDF2": "PyPDF2",
        "python-pptx": "pptx",
        "openpyxl": "openpyxl",
        "PyMuPDF": "fitz",
        "pdf2image": "pdf2image",
        "Pillow": "PIL",
        "pytesseract": "pytesseract",
        "pywin32": "win32com",
        "tqdm": "tqdm",
        "paddleocr": "paddleocr"  # 添加paddleocr依赖
    }
    
    print("\n检查Python包依赖...")
    
    # 检查每个依赖
    for package, module in dependencies.items():
        try:
            __import__(module)
            print(f"✓ {package} 已安装")
            installed_deps.append(package)
        except ImportError:
            print(f"× {package} 未安装")
            missing_deps.append(package)
    
    # 如果有缺失的依赖,尝试安装
    if missing_deps:
        print(f"\n发现 {len(missing_deps)} 个缺失的依赖,开始自动安装...")
        
        # 先升级pip
        try:
            print("\n升级pip...")
            subprocess.check_call([
                sys.executable, "-m", "pip", "install",
                "-i", "https://pypi.tuna.tsinghua.edu.cn/simple",
                "--upgrade", "pip"
            ])
        except:
            print("警告: pip升级失败,继续安装依赖...")
        
        # 安装缺失的依赖
        for package in missing_deps:
            if install_package(package):
                installed_deps.append(package)
            else:
                manual_deps.append(package)
        
        if manual_deps:
            print("\n以下依赖需要手动安装:")
            for package in manual_deps:
                print(f"pip install {package}")
            
            if "pdf2image" in manual_deps:
                print("\npdf2image还需要安装poppler:")
                print("1. 下载地址: https://github.com/oschwartz10612/poppler-windows/releases/")
                print("2. 解压到任意目录")
                print("3. 将bin目录添加到系统环境变量PATH")
        
        # 重新检查是否所有依赖都已安装
        missing_deps = []
        print("\n重新检查依赖...")
        for package, module in dependencies.items():
            if package not in manual_deps:  # 跳过需要手动安装的依赖
                try:
                    __import__(module)
                    print(f"✓ {package} 已安装")
                except ImportError:
                    print(f"× {package} 未安装")
                    missing_deps.append(package)
    
    # 检查Tesseract OCR
    print("\n检查Tesseract OCR...")
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        print("✓ Tesseract OCR 已安装")
        
        # 检查中文语言包
        if not os.path.exists(r'C:\Program Files\Tesseract-OCR\tessdata\chi_sim.traineddata'):
            print("\n! 警告：未找到中文语言包，OCR可能无法正确识别中文")
            print("请安装中文语言包:")
            print("1. 下载chi_sim.traineddata:")
            print("   https://github.com/tesseract-ocr/tessdata/raw/main/chi_sim.traineddata")
            print("2. 将文件放入 C:\\Program Files\\Tesseract-OCR\\tessdata 目录")
    except Exception as e:
        print("\n× Tesseract OCR 未安装或配置错误")
        print("请安装Tesseract OCR:")
        print("1. 下载地址: https://github.com/UB-Mannheim/tesseract/wiki")
        print("2. 选择最新版本的Windows安装包(如: tesseract-ocr-w64-setup-5.3.1.20230401.exe)")
        print("3. 运行安装程序,建议安装到默认目录: C:\\Program Files\\Tesseract-OCR")
        print("4. 安装时选择额外的语言包(Additional language data),确保选中Chinese simplified")
        print("5. 将安装目录添加到系统环境变量PATH")
        manual_deps.append("tesseract-ocr")
    
    return manual_deps  # 返回需要手动安装的依赖

class DocumentCompressor:
    # 添加类变量存储共享的OCR引擎
    _shared_ocr_engine = None
    _ocr_lock = mp.Lock()

    @classmethod
    def initialize_shared_ocr(cls):
        """在主进程中初始化OCR引擎"""
        if cls._shared_ocr_engine is None:
            try:
                # 设置日志级别
                import logging
                logging.getLogger("ppocr").setLevel(logging.ERROR)
                
                print("正在初始化OCR引擎...")
                from paddleocr import PaddleOCR
                
                # 确保模型文件夹存在
                model_dirs = [
                    'ch_PP-OCRv4_rec_infer',
                    'ch_PP-OCRv4_det_infer',
                    'ch_ppocr_mobile_v2.0_cls_infer'
                ]
                for dir_name in model_dirs:
                    if not os.path.exists(dir_name):
                        os.makedirs(dir_name, exist_ok=True)
                
                # 初始化OCR引擎
                cls._shared_ocr_engine = PaddleOCR(
                    use_angle_cls=True,
                    lang="ch",
                    use_gpu=False,
                    cpu_threads=4,
                    enable_mkldnn=True,
                    show_log=False,
                    rec_model_dir='ch_PP-OCRv4_rec_infer',
                    det_model_dir='ch_PP-OCRv4_det_infer',
                    cls_model_dir='ch_ppocr_mobile_v2.0_cls_infer',
                    check_install=False
                )
                print("OCR引擎初始化完成")
                
            except Exception as e:
                print(f"OCR引擎初始化失败: {str(e)}")
                cls._shared_ocr_engine = None

    def __init__(self):
        """初始化文档处理器"""
        self.input_dir = Path("raw")
        self.output_dir = Path("output")
        self.temp_dir = Path("temp")
        self.log_buffer = []
        
        # 创建必要的目录
        for dir_path in [self.input_dir, self.output_dir, self.temp_dir]:
            dir_path.mkdir(parents=True, exist_ok=True)
        
        # 支持的文件类型
        self.supported_extensions = {
            '.docx', '.doc',  # Word文档
            '.pdf',          # PDF文档
            '.pptx', '.ppt', # PPT文档
            '.xlsx', '.xls', # Excel文档
            '.txt',          # 纯文本
            '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'  # 图片文件
        }
        
        # 初始化统计信息
        self.stats = {
            'processed': 0,
            'failed': 0,
            'total_size_before': 0,
            'total_size_after': 0,
            'unsupported_files': [],
            'processing_times': {}
        }
        
        # 只在主进程中显示初始化信息
        if mp.current_process().name == 'MainProcess':
            print("\n文档压缩工具初始化完成")
            print(f"输入目录: {self.input_dir.absolute()}")
            print(f"输出目录: {self.output_dir.absolute()}")
            print(f"支持的文件类型: {', '.join(self.supported_extensions)}\n")

    def log(self, message):
        """添加日志到缓存"""
        self.log_buffer.append(message)

    def clear_log(self):
        """清空日志缓存"""
        self.log_buffer = []

    def print_log(self):
        """打印日志缓存"""
        for message in self.log_buffer:
            print(message)

    def check_files(self):
        """检查目录中的所有文件，返回支持和不支持的文件列表"""
        supported_files = []
        unsupported_files = []
        temp_files = []
        
        for file_path in self.input_dir.glob('*.*'):
            # 过滤临时文件
            if file_path.name.startswith('~$'):
                temp_files.append(file_path)
                continue
            
            # 检查文件扩展名
            if file_path.suffix.lower() in self.supported_extensions:
                # 检查文件是否可访问
                try:
                    with open(file_path, 'rb'):
                        supported_files.append(file_path)
                except PermissionError:
                    self.log(f"警告: 文件 {file_path.name} 被其他程序占用")
                    continue
            else:
                unsupported_files.append(file_path)
                self.stats['unsupported_files'].append(file_path.name)
        
        # 显示临时文件信息
        if temp_files:
            self.log("\n发现以下临时文件,将被跳过:")
            for temp_file in temp_files:
                self.log(f"- {temp_file.name}")
        
        return supported_files, unsupported_files

    def process_txt(self, file_path):
        """处理纯文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def process_doc(self, file_path):
        """处理.doc文件"""
        word = None
        try:
            # 在每个进程中创建新的COM对象
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(str(file_path.absolute()))
            text = doc.Content.Text
            doc.Close()
            return text
        except Exception as e:
            raise e
        finally:
            if word:
                try:
                    word.Quit()
                except:
                    pass

    def process_ppt(self, file_path):
        """处理.ppt文件"""
        powerpoint = None
        try:
            # 在每个进程中创建新的COM对象
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(str(file_path.absolute()))
            text = []
            
            for slide_num in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(slide_num)
                slide_text = [f"\n=== 第{slide_num}页 ===\n"]
                
                # 提取所有文字
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        text_frame = shape.TextFrame
                        if text_frame.HasText:
                            slide_text.append(text_frame.TextRange.Text)
                    
                    # 如果是图片，保存并OCR
                    if shape.Type == 13:  # msoPicture
                        try:
                            temp_image = self.temp_dir / f"slide_{slide_num}_image_{hash(str(shape.Id))}.png"
                            shape.Export(str(temp_image), "PNG")
                            
                            ocr_text = self.ocr_image(temp_image)
                            if ocr_text:
                                slide_text.append(f"[图片文字内容]\n{ocr_text}\n")
                            
                            # 删除临时图片
                            try:
                                os.remove(temp_image)
                            except:
                                pass
                        except Exception as e:
                            print(f"处理图片时出错: {str(e)}")
                
                text.extend(slide_text)
            
            presentation.Close()
            return '\n'.join(text)
        except Exception as e:
            raise e
        finally:
            if powerpoint:
                try:
                    powerpoint.Quit()
                except:
                    pass

    def process_xls(self, file_path):
        """处理.xls文件"""
        excel = None
        try:
            # 在每个进程中创建新的COM对象
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(str(file_path.absolute()))
            text = []
            
            # 统计工作表信息
            sheet_count = wb.Sheets.Count
            self.log(f"工作簿统计: {sheet_count} 个工作表")
            
            # 处理每个工作表
            total_rows = 0
            non_empty_rows = 0
            for sheet in wb.Sheets:
                self.log(f"处理工作表: {sheet.Name}")
                used_range = sheet.UsedRange
                sheet_rows = used_range.Rows.Count
                total_rows += sheet_rows
                
                sheet_non_empty = 0
                for row in used_range.Rows:
                    row_text = []
                    for cell in row.Cells:
                        if cell.Text:
                            row_text.append(str(cell.Text).strip())
                    if row_text:
                        text.append(' | '.join(row_text))
                        non_empty_rows += 1
                        sheet_non_empty += 1
                self.log(f"提取了 {sheet_non_empty} 行非空数据")
            
            wb.Close()
            return '\n'.join(text)
        except Exception as e:
            raise e
        finally:
            if excel:
                try:
                    excel.Quit()
                except:
                    pass

    def compress_text(self, text):
        """压缩文本内容,去除重复和无用信息"""
        if not text:
            return text
            
        # 删除多余空白,但保留段落分隔
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        lines = text.split('\n')
        unique_lines = []
        seen = set()
        
        # 页眉页脚模式识别
        header_footer_patterns = set()
        for line in lines[:5] + lines[-5:]:  # 只检查开头和结尾的少量行
            line = line.strip()
            if len(line) > 20:  # 增加长度阈值,避免误删
                header_footer_patterns.add(line)
        
        # 处理每一行
        for line in lines:
            line = line.strip()
            if not line:  # 保留空行作为段落分隔
                unique_lines.append('')
                continue
                
            # 跳过页眉页脚(完全匹配)
            if line in header_footer_patterns:
                continue
                
            # 跳过过短的行
            if len(line) < 2:
                continue
                
            # 相似度去重
            skip_line = False
            for existing_line in list(seen)[-5:]:  # 只检查最近的5行
                if self._is_similar(line, existing_line):
                    skip_line = True
                    break
            
            if not skip_line:
                seen.add(line)
                unique_lines.append(line)
        
        # 合并处理后的文本,保留段落格式
        compressed = '\n'.join(unique_lines)
        
        # 如果压缩后为空,返回原始文本
        if not compressed.strip():
            return text
            
        # 如果还是太长,进行截断
        max_length = 1000000  # 增加长度限制
        if len(compressed) > max_length:
            half_length = max_length // 2
            compressed = compressed[:half_length] + "\n...(内容已截断)...\n" + compressed[-half_length:]
        
        return compressed

    def _is_similar(self, str1, str2):
        """检查两个字符串是否相似"""
        # 如果长度相差太多,直接认为不相似
        if abs(len(str1) - len(str2)) > min(len(str1), len(str2)) * 0.5:  # 增加差异容忍度
            return False
            
        # 计算编辑距离
        if len(str1) > len(str2):
            str1, str2 = str2, str1
            
        distances = range(len(str1) + 1)
        for i2, c2 in enumerate(str2):
            distances_ = [i2+1]
            for i1, c1 in enumerate(str1):
                if c1 == c2:
                    distances_.append(distances[i1])
                else:
                    distances_.append(1 + min((distances[i1], distances[i1 + 1], distances_[-1])))
            distances = distances_
            
        # 降低相似度阈值
        similarity = 1 - distances[-1] / max(len(str1), len(str2))
        return similarity > 0.9  # 提高到90%相似度才算重复

    def process_docx(self, file_path):
        """处理.docx文件"""
        self.log(f"处理Word文件: {file_path.name}")
        doc = Document(file_path)
        text = []
        
        # 统计文档信息
        para_count = len(doc.paragraphs)
        table_count = len(doc.tables)
        self.log(f"文档统计: {para_count}段落, {table_count}表格")
        
        # 提取段落
        para_extracted = 0
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text.strip())
                para_extracted += 1
        self.log(f"提取了{para_extracted}个非空段落")
                
        # 提取表格
        table_rows = 0
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text.append(row_text)
                    table_rows += 1
        self.log(f"提取了{table_rows}行表格数据")
        
        # 检查提取的文本
        raw_text = '\n'.join(text)
        raw_size = len(raw_text)
        self.log(f"提取的原始文本: {raw_size}字符")
        
        if raw_size < 100:  # 如果提取的文本太少
            self.log("警告：提取的文本内容过少,尝试使用COM对象")
            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                try:
                    doc = word.Documents.Open(str(file_path.absolute()))
                    com_text = doc.Content.Text
                    doc.Close()
                    word.Quit()
                    
                    com_size = len(com_text)
                    if com_size > raw_size:
                        self.log(f"使用COM对象提取: {com_size}字符")
                        raw_text = com_text
                    else:
                        self.log("保留python-docx提取的文本")
                        
                except Exception as e:
                    self.log(f"COM对象处理失败: {str(e)}")
                    if word:
                        word.Quit()
            except Exception as e:
                self.log(f"备选方法失败: {str(e)}")
        
        # 确保有提取到内容
        if not raw_text.strip():
            raise Exception("无法从文档中提取文本")
            
        return raw_text

    def convert_pdf_page_to_images(self, file_path, page_num):
        """将PDF页面转换为图片列表"""
        images = []
        try:
            # 方法1: 使用PyMuPDF
            doc = fitz.open(file_path)
            page = doc[page_num - 1]
            # 提高缩放比例到1.5
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
            doc.close()
        except Exception as e:
            print(f"PyMuPDF转换失败: {str(e)}")
            try:
                # 方法2: 使用pdf2image,提高DPI到200
                pages = convert_from_path(file_path, first_page=page_num, last_page=page_num, dpi=200)
                images.extend(pages)
            except Exception as e:
                print(f"pdf2image转换失败: {str(e)}")
        return images

    def preprocess_image(self, image):
        """图像预处理优化"""
        # 将PIL Image转换为OpenCV格式
        if isinstance(image, Image.Image):
            image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        
        # 提高分辨率上限
        max_dimension = 2880
        
        # 计算缩放比例
        height, width = image.shape[:2]
        scale = max_dimension / max(height, width)
        if scale > 1:
            image = cv2.resize(image, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
        
        # 自适应直方图均衡化
        if len(image.shape) == 3:
            lab = cv2.cvtColor(image, cv2.COLOR_BGR2LAB)
            l, a, b = cv2.split(lab)
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
            l = clahe.apply(l)
            lab = cv2.merge((l,a,b))
            image = cv2.cvtColor(lab, cv2.COLOR_LAB2BGR)
        
        # 锐化
        kernel = np.array([[-1,-1,-1], [-1,9,-1], [-1,-1,-1]])
        image = cv2.filter2D(image, -1, kernel)
        
        # 降噪
        image = cv2.fastNlMeansDenoisingColored(image, None, 10, 10, 7, 21)
            
        return image

    def split_image(self, image, max_height=1500):  # 降低最大高度
        """将长图片分割成多个部分"""
        try:
            if image.height <= max_height:
                return [image]
            
            parts = []
            current_y = 0
            
            while current_y < image.height:
                # 计算当前部分的高度
                part_height = min(max_height, image.height - current_y)
                
                # 添加重叠区域
                overlap = 50  # 降低重叠区域到50像素
                if current_y > 0:
                    current_y -= overlap
                
                # 裁剪图片
                part = image.crop((0, current_y, image.width, current_y + part_height))
                parts.append(part)
                
                current_y += part_height
                
                # 添加进度提示
                self.log(f"已分割图片: {len(parts)}/{int(image.height/max_height + 1)}")
            
            return parts
        except Exception as e:
            self.log(f"图片分割失败: {str(e)}")
            return [image]

    def init_ocr(self):
        """获取OCR引擎实例"""
        with self._ocr_lock:
            if self._shared_ocr_engine is None:
                self.initialize_shared_ocr()
            return self._shared_ocr_engine

    def ocr_image(self, image_path_or_image, ocr_engine=None, max_retries=3):
        """OCR处理添加重试机制"""
        for attempt in range(max_retries):
            try:
                # 加载图像
                if isinstance(image_path_or_image, str) or isinstance(image_path_or_image, Path):
                    image = Image.open(image_path_or_image)
                else:
                    image = image_path_or_image
                
                # 图像预处理
                processed_image = self.preprocess_image(image)
                
                # OCR识别
                result = ocr_engine.ocr(processed_image)
                
                # 提取文本
                text_parts = []
                if result is not None:
                    for line in result:
                        if line is not None and len(line) >= 2:
                            # PaddleOCR返回格式: [[box坐标], (文本内容, 置信度)]
                            # 提取文本内容
                            text_content = line[1][0]
                            if text_content and isinstance(text_content, str):
                                text_parts.append(text_content)
                
                # 合并所有文本
                text = '\n'.join(text_parts) if text_parts else ''
                
                # 后处理
                text = self.postprocess_text(text)
                
                return text
                
            except Exception as e:
                if attempt == max_retries - 1:
                    print(f"OCR处理失败: {str(e)}")
                    return ""
                time.sleep(1)  # 重试前等待

    def process_pdf(self, file_path):
        """处理PDF文件,支持文本型和图片型PDF"""
        self.log(f"处理PDF文件: {file_path.name}")
        text_parts = []
        
        try:
            # 首先尝试使用 PyMuPDF 处理
            doc = fitz.open(file_path)
            total_pages = len(doc)
            self.log(f"总页数: {total_pages}")
            
            for page_num in range(total_pages):
                try:
                    page = doc[page_num]
                    # 先尝试提取文本
                    page_text = page.get_text()
                    
                    # 如果文本太少,尝试OCR
                    if len(page_text.strip()) < 100:
                        self.log(f"第{page_num + 1}页文本内容过少,尝试OCR处理...")
                        # 将页面转换为图片,使用较小的缩放比例
                        pix = page.get_pixmap(matrix=fitz.Matrix(1.0, 1.0))  # 降低缩放比例
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        # OCR处理
                        ocr_text = self.ocr_image(img)
                        if ocr_text:
                            page_text = ocr_text
                    
                    if page_text.strip():
                        text_parts.append(page_text)
                        self.log(f"第{page_num + 1}页提取了{len(page_text)}字符")
                    else:
                        self.log(f"第{page_num + 1}页未能提取到内容")
                    
                except Exception as e:
                    self.log(f"警告：处理第{page_num + 1}页时出错: {str(e)}")
                    continue
                
            doc.close()
            
        except Exception as e:
            self.log(f"PyMuPDF处理失败: {str(e)}")
            # 如果 PyMuPDF 失败,尝试使用 pdf2image
            try:
                self.log("尝试使用 pdf2image 处理...")
                images = convert_from_path(file_path, dpi=150)  # 降低DPI
                for i, image in enumerate(images):
                    try:
                        ocr_text = self.ocr_image(image)
                        if ocr_text:
                            text_parts.append(ocr_text)
                            self.log(f"第{i + 1}页OCR提取了{len(ocr_text)}字符")
                    except Exception as e:
                        self.log(f"警告：第{i + 1}页OCR处理失败: {str(e)}")
                        continue
            except Exception as e:
                self.log(f"pdf2image处理失败: {str(e)}")
                # 最后尝试使用 PyPDF2
                try:
                    self.log("尝试使用 PyPDF2 处理...")
                    reader = PdfReader(file_path)
                    for page_num, page in enumerate(reader.pages):
                        try:
                            page_text = page.extract_text()
                            if page_text.strip():
                                text_parts.append(page_text)
                                self.log(f"第{page_num + 1}页提取了{len(page_text)}字符")
                        except Exception as e:
                            self.log(f"警告：第{page_num + 1}页文本提取失败: {str(e)}")
                            continue
                except Exception as e:
                    self.log(f"PyPDF2处理失败: {str(e)}")
        
        # 合并所有文本
        all_text = '\n'.join(text_parts)
        if not all_text.strip():
            raise Exception("无法从PDF中提取文本")
        
        self.log(f"总计提取了 {len(all_text)} 字符")
        return all_text

    def process_pptx(self, file_path):
        """处理.pptx文件"""
        self.log(f"处理PPT文件: {file_path.name}")
        prs = Presentation(file_path)
        text = []
        
        total_slides = len(prs.slides)
        self.log(f"PPT统计: {total_slides}页")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                slide_text = [f"\n=== 第{slide_num}页 ===\n"]
                shape_count = len(slide.shapes)
                
                # 提取所有文字
                text_extracted = 0
                for shape in slide.shapes:
                    try:
                        shape_text = self.extract_text_from_shape(shape)
                        if shape_text.strip():
                            slide_text.append(shape_text)
                            text_extracted += 1
                    except Exception as e:
                        self.log(f"警告：提取形状文字时出错: {str(e)}")
                
                # 提取并OCR图片
                images_processed = 0
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "image"):
                            image_path = self.extract_image_from_shape(shape, slide_num)
                            if image_path:
                                ocr_text = self.ocr_image(image_path)
                                if ocr_text:
                                    slide_text.append(f"[图片文字内容]\n{ocr_text}\n")
                                    images_processed += 1
                                
                                # 删除临时图片
                                try:
                                    os.remove(image_path)
                                except:
                                    pass
                    except Exception as e:
                        self.log(f"警告：处理图片时出错: {str(e)}")
                
                text.extend(slide_text)
                self.log(f"第{slide_num}页: {text_extracted}个文本, {images_processed}张图片")
                
            except Exception as e:
                self.log(f"警告：处理第{slide_num}页时出错: {str(e)}")
                continue
        
        # 合并所有文本
        all_text = '\n'.join(text)
        if not all_text.strip():
            raise Exception("无法从PPT中提取文本")
        
        self.log(f"总计提取了{len(all_text)}字符")
        return all_text

    def process_xlsx(self, file_path):
        """处理.xlsx文件"""
        self.log(f"处理Excel文件: {file_path.name}")
        wb = load_workbook(file_path)
        text = []
        
        # 统计工作表信息
        sheet_count = len(wb.worksheets)
        self.log(f"工作簿统计: {sheet_count}个工作表")
        
        # 处理每个工作表
        total_rows = 0
        non_empty_rows = 0
        for sheet in wb.worksheets:
            sheet_rows = 0
            for row in sheet.rows:
                total_rows += 1
                row_text = []
                for cell in row:
                    if cell.value:
                        row_text.append(str(cell.value).strip())
                if row_text:
                    text.append(' | '.join(row_text))
                    non_empty_rows += 1
                    sheet_rows += 1
            self.log(f"工作表[{sheet.title}]: {sheet_rows}行")
        
        self.log(f"总计处理: {total_rows}行中提取{non_empty_rows}行")
        
        # 检查提取的文本
        raw_text = '\n'.join(text)
        raw_size = len(raw_text)
        self.log(f"提取的原始文本: {raw_size}字符")
        
        # 如果提取的内容太少,尝试使用COM对象
        if raw_size < 100:
            self.log("警告：提取的文本内容过少,尝试使用COM对象")
            try:
                com_text = self.process_xls(file_path)
                com_size = len(com_text)
                
                # 如果COM提取的内容更多,使用COM结果
                if com_size > raw_size:
                    self.log(f"使用COM对象提取: {com_size}字符")
                    raw_text = com_text
                else:
                    self.log("保留openpyxl提取的文本")
            except Exception as e:
                self.log(f"COM对象处理失败: {str(e)}")
        
        # 确保有提取到内容
        if not raw_text.strip():
            raise Exception("无法从Excel文件中提取文本")
            
        return raw_text

    def extract_text_from_shape(self, shape):
        """从PPT形状中提取文字"""
        text = ""
        
        try:
            # 文本框
            if hasattr(shape, "text"):
                text += shape.text + "\n"
                
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                        
            # SmartArt
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
                        
            # 组合形状
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    text += self.extract_text_from_shape(sub_shape)
        except Exception as e:
            self.log(f"警告：提取形状文字时出错: {str(e)}")
            
        return text

    def extract_image_from_shape(self, shape, slide_num):
        """从PPT形状中提取图片"""
        try:
            if hasattr(shape, "image"):
                self.log("读取图片数据...")
                image_bytes = shape.image.blob
                self.log("转换图片格式...")
                image = Image.open(io.BytesIO(image_bytes))
                
                # 保存图片到临时目录
                image_path = self.temp_dir / f"slide_{slide_num}_image_{hash(str(image_bytes))}.png"
                self.log(f"保存临时图片: {image_path}")
                image.save(image_path)
                return image_path
        except Exception as e:
            self.log(f"警告：提取图片时出错: {str(e)}")
        return None

    def process_image(self, file_path, ocr_engine=None):
        """处理图片文件"""
        try:
            self.log(f"处理图片文件: {file_path.name}")
            
            # 读取图片
            image = Image.open(file_path)
            
            # 获取原始图片信息
            original_size = image.size
            self.log(f"原始图片尺寸: {original_size}")
            
            # 如果是长截图,直接分割处理
            if original_size[1] > original_size[0] * 3:  # 高度超过宽度的3倍
                self.log("检测到长截图,开始分段处理...")
                # 计算需要分割的段数
                segment_height = 1500  # 每段高度
                total_segments = (original_size[1] + segment_height - 1) // segment_height
                self.log(f"图片将分割为 {total_segments} 段处理")
                
                # 分段处理
                all_text = []
                processed_segments = 0
                error_count = 0  # 添加错误计数
                
                for i in range(total_segments):
                    try:
                        # 计算当前段的区域
                        start_y = i * segment_height
                        end_y = min((i + 1) * segment_height, original_size[1])
                        
                        # 裁剪当前段
                        segment = image.crop((0, start_y, original_size[0], end_y))
                        
                        # OCR处理当前段
                        segment_text = self.ocr_image(segment, ocr_engine)
                        if segment_text.strip():
                            all_text.append(segment_text)
                        
                        # 释放内存
                        del segment
                        
                        # 更新进度
                        processed_segments += 1
                        print(f"\r图片分段处理: {processed_segments}/{total_segments} ({processed_segments/total_segments*100:.1f}%)", end="", flush=True)
                        
                    except Exception as e:
                        error_count += 1
                        if error_count <= 3:  # 只显示前3个错误
                            print(f"\nOCR处理失败: {str(e)}")
                        continue
                
                print()  # 换行
                
                # 合并所有文本
                text = '\n'.join(all_text)
                
            else:
                # 普通图片直接处理
                text = self.ocr_image(image, ocr_engine)
            
            if not text.strip():
                raise Exception("OCR未能识别出文本")
            
            # 压缩文本
            self.log("压缩识别出的文本...")
            compressed_text = self.compress_text(text)
            
            # 保存结果
            output_file = self.output_dir / f"{file_path.stem}_compressed.txt"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(compressed_text)
            
            # 更新统计信息
            self.stats['total_size_before'] += file_path.stat().st_size
            self.stats['total_size_after'] += output_file.stat().st_size
            self.stats['processed'] += 1
            
            # 显示处理结果
            print(f"✓ {file_path.name} 处理成功")
            print(f"  - 原始大小: {file_path.stat().st_size / 1024:.2f} KB")
            print(f"  - 压缩大小: {output_file.stat().st_size / 1024:.2f} KB")
            print(f"  - 压缩比: {output_file.stat().st_size / file_path.stat().st_size:.2%}")
            
            return True
            
        except Exception as e:
            self.log(f"图片处理失败: {str(e)}")
            print(f"\n× {file_path.name} 处理失败")
            print(f"错误原因: {str(e)}")
            print("\n详细日志:")
            self.print_log()
            self.clear_log()
            return False

    def clear_temp_files(self):
        """清理temp目录中的所有临时文件"""
        try:
            for file in self.temp_dir.glob("*"):
                if file.is_file():
                    try:
                        file.unlink()
                    except Exception as e:
                        self.log(f"警告：删除临时文件 {file.name} 失败: {str(e)}")
        except Exception as e:
            self.log(f"清理临时文件时出错: {str(e)}")

    def process_file(self, file_path):
        try:
            file_path = Path(file_path)
            if file_path.suffix.lower() not in self.supported_extensions:
                return
            
            self.log(f"读取文件大小: {file_path.stat().st_size / 1024:.2f} KB")
            self.stats['total_size_before'] += file_path.stat().st_size
            
            # 处理文件
            ext = file_path.suffix.lower()
            self.log(f"文件类型: {ext}")
            
            # 处理图片文件
            if ext in {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'}:
                self.process_image(file_path)
                return
                
            # 处理其他文件类型
            if ext == '.docx':
                text = self.process_docx(file_path)
            elif ext == '.doc':
                text = self.process_doc(file_path)
            elif ext == '.pdf':
                text = self.process_pdf(file_path)
            elif ext == '.pptx':
                text = self.process_pptx(file_path)
            elif ext == '.ppt':
                text = self.process_ppt(file_path)
            elif ext == '.xlsx':
                text = self.process_xlsx(file_path)
            elif ext == '.xls':
                text = self.process_xls(file_path)
            elif ext == '.txt':
                text = self.process_txt(file_path)
            
            # 检查提取的文本
            if not text:
                raise Exception("文本提取失败")
                
            # 显示原始文本信息
            raw_lines = text.split('\n')
            self.log(f"提取的文本统计:")
            self.log(f"- 字符数: {len(text)}")
            self.log(f"- 行数: {len(raw_lines)}")
            self.log(f"- 非空行: {len([l for l in raw_lines if l.strip()])}")
            
            self.log("压缩文本...")
            compressed_text = self.compress_text(text)
            
            # 显示压缩后文本信息
            comp_lines = compressed_text.split('\n')
            self.log(f"压缩后文本统计:")
            self.log(f"- 字符数: {len(compressed_text)}")
            self.log(f"- 行数: {len(comp_lines)}")
            self.log(f"- 非空行: {len([l for l in comp_lines if l.strip()])}")
            self.log(f"- 压缩比: {len(compressed_text) / len(text):.2%}")
            
            # 检查压缩结果
            if len(compressed_text.strip()) == 0:
                self.log("警告：压缩后文本为空，使用原始文本")
                compressed_text = text
            
            # 保存文件
            output_file = self.output_dir / f"{file_path.stem}_compressed.txt"
            self.log(f"保存到: {output_file}")
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(compressed_text)
            
            self.stats['total_size_after'] += output_file.stat().st_size
            self.stats['processed'] += 1
            
            # 处理成功,只显示简要信息
            print(f"✓ {file_path.name} 处理成功")
            print(f"  - 原始大小: {file_path.stat().st_size / 1024:.2f} KB")
            print(f"  - 压缩大小: {output_file.stat().st_size / 1024:.2f} KB")
            print(f"  - 压缩比: {output_file.stat().st_size / file_path.stat().st_size:.2%}")
            
            # 清空日志缓存
            self.clear_log()
            
        except Exception as e:
            # 处理失败,显示详细日志
            print(f"\n× {file_path.name} 处理失败")
            print(f"错误原因: {str(e)}")
            print("\n详细日志:")
            self.print_log()
            self.clear_log()
            raise e
        finally:
            # 清理临时文件
            self.clear_temp_files()

    def process_all_files(self):
        """处理所有文件"""
        import signal
        
        def signal_handler(signum, frame):
            print('\n\n收到中断信号，正在停止处理...')
            if 'pool' in locals():
                pool.terminate()
                pool.join()
            self.print_stats()
            self.clear_temp_files()
            sys.exit(0)
        
        # 注册信号处理器
        signal.signal(signal.SIGINT, signal_handler)
        
        # 检查目录
        if not self.input_dir.exists():
            print(f"创建输入目录: {self.input_dir}")
            self.input_dir.mkdir(parents=True)
            print("请将要处理的文件放入 raw 目录中，然后重新运行程序")
            return

        if not self.output_dir.exists():
            self.output_dir.mkdir(parents=True)

        # 检查文件
        supported_files, unsupported_files = self.check_files()

        if not supported_files:
            print("\n未找到可处理的文件！")
            print("请将文件放入 raw 目录后重试")
            return

        print(f"\n找到 {len(supported_files)} 个文件待处理")
        
        # 获取最优进程数
        num_processes = get_optimal_processes()
        print(f"将使用 {num_processes} 个进程并行处理")
        
        try:
            # 使用上下文管理器创建进程池
            with Pool(processes=num_processes) as pool:
                results = []
                with tqdm(total=len(supported_files), desc="总进度", position=0) as pbar:
                    try:
                        for result in pool.imap(self._process_file_wrapper, supported_files):
                            results.append(result)
                            pbar.update(1)
                            
                            # 只更新统计信息,不输出详细日志
                            if result['success']:
                                self.stats['processed'] += 1
                                self.stats['total_size_before'] += result['size_before']
                                self.stats['total_size_after'] += result['size_after']
                            else:
                                self.stats['failed'] += 1
                            
                    except KeyboardInterrupt:
                        print('\n\n收到中断信号，正在停止处理...')
                        pool.terminate()
                        raise
                    
        except KeyboardInterrupt:
            pass
        finally:
            # 清理并显示统计
            self.clear_temp_files()
            self.print_stats()

    def _process_file_wrapper(self, file_path):
        """处理单个文件的包装函数，用于多进程处理"""
        result = {
            'success': False,
            'size_before': 0,
            'size_after': 0,
            'file_name': file_path.name
        }
        
        try:
            # 使用共享的OCR引擎
            ocr_engine = self.init_ocr()
            if ocr_engine is None:
                raise Exception("OCR引擎初始化失败")
            
            # 检查文件是否被占用
            try:
                with open(file_path, 'rb'):
                    pass
            except PermissionError:
                print(f"× {file_path.name} - 文件被占用")
                return result
            
            # 获取原始文件大小
            result['size_before'] = file_path.stat().st_size
            
            # 处理文件
            if file_path.suffix.lower() in {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'}:
                success = self.process_image(file_path, ocr_engine)
            else:
                # 处理其他类型文件
                ext = file_path.suffix.lower()
                if ext == '.docx':
                    text = self.process_docx(file_path)
                elif ext == '.doc':
                    text = self.process_doc(file_path)
                elif ext == '.pdf':
                    text = self.process_pdf(file_path)
                elif ext == '.pptx':
                    text = self.process_pptx(file_path)
                elif ext == '.ppt':
                    text = self.process_ppt(file_path)
                elif ext == '.xlsx':
                    text = self.process_xlsx(file_path)
                elif ext == '.xls':
                    text = self.process_xls(file_path)
                elif ext == '.txt':
                    text = self.process_txt(file_path)
                
                if not text:
                    raise Exception("文本提取失败")
                
                # 压缩文本
                compressed_text = self.compress_text(text)
                if len(compressed_text.strip()) == 0:
                    compressed_text = text
                
                # 保存文件
                output_file = self.output_dir / f"{file_path.stem}_compressed.txt"
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(compressed_text)
                
                success = True
                result['size_after'] = output_file.stat().st_size
            
            if success:
                result['success'] = True
                print(f"✓ {file_path.name} 处理成功")
                print(f"  - 原始大小: {result['size_before'] / 1024:.2f} KB")
                print(f"  - 压缩大小: {result['size_after'] / 1024:.2f} KB")
                print(f"  - 压缩比: {result['size_after'] / result['size_before']:.2%}")
            
        except Exception as e:
            print(f"\n× {file_path.name} 处理失败")
            print(f"错误原因: {str(e)}")
        finally:
            # 清理临时文件
            self.clear_temp_files()
        
        return result
    
    def print_stats(self):
        """打印处理统计信息"""
        print("\n处理统计:")
        print(f"✓ 成功处理: {self.stats['processed']} 个文件")
        if self.stats['failed'] > 0:
            print(f"× 处理失败: {self.stats['failed']} 个文件")
        if self.stats['unsupported_files']:
            print(f"- 不支持的文件: {len(self.stats['unsupported_files'])} 个")
        print(f"- 原始大小: {self.stats['total_size_before'] / 1024:.2f} KB")
        print(f"- 压缩大小: {self.stats['total_size_after'] / 1024:.2f} KB")
        if self.stats['total_size_before'] > 0:
            ratio = self.stats['total_size_after'] / self.stats['total_size_before']
            print(f"- 总压缩比: {ratio:.2%}")
            print(f"- 节省空间: {(self.stats['total_size_before'] - self.stats['total_size_after']) / 1024:.2f} KB")

    def postprocess_text(self, text):
        """文本后处理优化"""
        import re
        
        # 修正常见错误
        corrections = {
            '糖': '需',
            '斋': '需',
            # 添加更多常见错误修正
        }
        
        # 应用修正
        for wrong, right in corrections.items():
            text = text.replace(wrong, right)
        
        # 修正标点符号
        text = re.sub(r'[，,]+', '，', text)  # 统一中文逗号
        text = re.sub(r'[。.]+', '。', text)  # 统一中文句号
        text = re.sub(r'[、]+', '、', text)  # 处理顿号
        
        # 修正空格
        text = re.sub(r'\s+', ' ', text)  # 合并多个空格
        text = re.sub(r'([。，！？])([^"\'])', r'\1\n\2', text)  # 在句末添加换行
        
        # 修正特殊符号
        text = text.replace('0', '•').replace('O', '•')  # 修正圆点符号
        
        return text

def main():
    print("文档压缩工具 v1.0")
    print("=" * 50)
    
    # 检查依赖
    manual_deps = check_dependencies()
    if manual_deps:
        print(f"\n以下依赖需要手动安装,请按照上述说明进行安装:")
        for dep in manual_deps:
            print(f"- {dep}")
        print("\n安装完成后重新运行程序")
        input("\n按回车键退出...")
        sys.exit(1)
    
    try:
        # 在主进程中初始化OCR引擎
        DocumentCompressor.initialize_shared_ocr()
        
        compressor = DocumentCompressor()
        print("\n支持格式:", ", ".join(sorted(compressor.supported_extensions)))
        print("\n使用说明:")
        print("1. 将需要处理的文件放在 raw 目录下")
        print("2. 运行程序进行处理")
        print("3. 处理后的文件保存在 output 目录下")
        print("4. 按 Ctrl+C 可随时中断处理")
        print("=" * 50)
        
        compressor.process_all_files()
        
    except KeyboardInterrupt:
        print("\n程序已中断")
    except Exception as e:
        print(f"\n程序出错: {str(e)}")
    finally:
        print("\n处理完成!")

if __name__ == "__main__":
    main() 